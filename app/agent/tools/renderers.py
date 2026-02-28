"""
ARIA — Renderers
Converts the final report dict into markdown, html, pdf, pptx.
"""

from __future__ import annotations

import json
import os
from datetime import datetime
from typing import Any


OUTPUT_DIR = "outputs"


def _ensure_output_dir() -> str:
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    return OUTPUT_DIR


def _ts() -> str:
    return datetime.utcnow().strftime("%Y%m%d_%H%M%S")


# ──────────────────────────────────────────────────────────────
# MARKDOWN
# ──────────────────────────────────────────────────────────────

def render_markdown(report: dict) -> str:
    """Convert report dict to a markdown string."""
    overview  = report.get("1_overview", {})
    summary   = report.get("2_data_summary", {})
    triz      = report.get("3_triz_analysis", {})
    findings  = report.get("4_key_findings", [])
    recs      = report.get("5_recommendations", [])
    conf      = report.get("6_confidence", {})

    lines = [
        f"# ARIA Activity Report",
        f"**Domain:** {overview.get('domain', 'N/A')}  ",
        f"**Period:** {overview.get('reporting_period', 'N/A')}  ",
        f"**Generated:** {overview.get('generated_at', 'N/A')}  ",
        f"**KPIs:** {', '.join(overview.get('kpis', []))}",
        "",
        "---",
        "",
        "## 1. Data Summary",
        f"- Records consolidated: {summary.get('consolidated_dataset', {}).get('record_count', 'N/A')}",
        f"- RAG chunks used: {summary.get('rag_chunks_used', 0)}",
        "",
        "## 2. TRIZ Analysis",
    ]

    for c in triz.get("contradictions", []):
        lines.append(f"- **{c.get('type','').upper()} contradiction**: "
                     f"{c.get('improving_parameter')} ↑ vs {c.get('degrading_parameter')} ↓  "
                     f"— {c.get('description','')}")

    lines += [
        "",
        f"**Ideal Final Result:** {triz.get('ideal_final_result', 'N/A')}",
        "",
        "**TRIZ Principles Applied:**",
    ]
    for p in triz.get("triz_principles_applied", []):
        lines.append(f"- #{p.get('principle_number')} {p.get('name')}: {p.get('application')}")

    lines += ["", "**Root Causes:**"]
    for rc in triz.get("root_causes", []):
        lines.append(f"- {rc}")

    lines += ["", "## 3. Key Findings"]
    for i, f in enumerate(findings, 1):
        lines.append(f"{i}. {f}")

    lines += ["", "## 4. Recommendations"]
    for r in recs:
        lines.append(
            f"- **[{r.get('priority','?')}]** {r.get('action')} "
            f"| Owner: {r.get('owner','?')} | Timeline: {r.get('timeline','?')}"
        )

    degraded = conf.get("degraded", False)
    lines += [
        "",
        "---",
        f"## 5. Confidence Score: {conf.get('percent', 'N/A')}",
        "⚠️ *Degraded report — max retries reached with low confidence.*" if degraded else "",
    ]

    return "\n".join(lines)


# ──────────────────────────────────────────────────────────────
# HTML
# ──────────────────────────────────────────────────────────────

def render_html(report: dict, chart_paths: list[str] | None = None) -> str:
    """Convert report dict to an HTML string with embedded chart images."""
    md = render_markdown(report)
    charts_html = ""
    for path in (chart_paths or []):
        charts_html += f'<img src="{path}" style="max-width:100%;margin:12px 0;" />\n'

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <title>ARIA Report — {report.get("1_overview", {}).get("domain", "")}</title>
  <style>
    body {{ font-family: Arial, sans-serif; max-width: 900px; margin: 40px auto; padding: 0 20px; }}
    h1 {{ color: #2c3e50; }} h2 {{ color: #34495e; border-bottom: 1px solid #eee; padding-bottom: 4px; }}
    pre {{ background: #f4f4f4; padding: 12px; border-radius: 4px; overflow-x: auto; }}
  </style>
</head>
<body>
<pre>{md}</pre>
{charts_html}
</body>
</html>"""


# ──────────────────────────────────────────────────────────────
# PDF
# ──────────────────────────────────────────────────────────────

def render_pdf(report: dict, chart_paths: list[str] | None = None) -> str:
    """
    Generate a PDF from the report and return its file path.
    Requires: pip install weasyprint
    """
    try:
        from weasyprint import HTML as WeasyprintHTML
    except ImportError:
        raise ImportError("weasyprint required: pip install weasyprint")

    html_content = render_html(report, chart_paths)
    out_dir  = _ensure_output_dir()
    out_path = os.path.join(out_dir, f"aria_report_{_ts()}.pdf")
    WeasyprintHTML(string=html_content).write_pdf(out_path)
    return out_path


# ──────────────────────────────────────────────────────────────
# PPTX
# ──────────────────────────────────────────────────────────────

def render_pptx(report: dict, chart_paths: list[str] | None = None) -> str:
    """
    Generate a PowerPoint from the report and return its file path.
    Requires: pip install python-pptx
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("python-pptx required: pip install python-pptx")

    overview = report.get("1_overview", {})
    triz     = report.get("3_triz_analysis", {})
    findings = report.get("4_key_findings", [])
    recs     = report.get("5_recommendations", [])
    conf     = report.get("6_confidence", {})

    prs = Presentation()
    blank = prs.slide_layouts[1]  # title + content

    def add_slide(title: str, content: str):
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    # Slide 1 — Overview
    add_slide(
        f"ARIA Report — {overview.get('domain', 'N/A')}",
        f"Period: {overview.get('reporting_period', 'N/A')}\n"
        f"KPIs: {', '.join(overview.get('kpis', []))}\n"
        f"Generated: {overview.get('generated_at', 'N/A')}"
    )

    # Slide 2 — TRIZ Contradictions
    contradictions = "\n".join(
        f"• {c.get('type','').upper()}: {c.get('improving_parameter')} ↑ vs {c.get('degrading_parameter')} ↓"
        for c in triz.get("contradictions", [])
    )
    add_slide("TRIZ — Contradictions", contradictions or "No contradictions identified.")

    # Slide 3 — IFR + Principles
    principles = "\n".join(
        f"• #{p.get('principle_number')} {p.get('name')}: {p.get('application')}"
        for p in triz.get("triz_principles_applied", [])
    )
    add_slide(
        "TRIZ — IFR & Principles",
        f"IFR: {triz.get('ideal_final_result', 'N/A')}\n\n{principles}"
    )

    # Slide 4 — Key Findings
    add_slide("Key Findings", "\n".join(f"{i+1}. {f}" for i, f in enumerate(findings)))

    # Slide 5 — Recommendations
    rec_text = "\n".join(
        f"• [{r.get('priority','?')}] {r.get('action')} | {r.get('owner','?')} | {r.get('timeline','?')}"
        for r in recs
    )
    add_slide("Recommendations", rec_text or "No recommendations.")

    # Slide 6 — Confidence
    degraded_note = "\n⚠️ Degraded report." if conf.get("degraded") else ""
    add_slide("Confidence Score", f"{conf.get('percent', 'N/A')}{degraded_note}")

    # Embed charts if available
    for path in (chart_paths or []):
        if os.path.exists(path):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            slide.shapes.add_picture(path, Inches(0.5), Inches(0.5), Inches(9))

    out_dir  = _ensure_output_dir()
    out_path = os.path.join(out_dir, f"aria_report_{_ts()}.pptx")
    prs.save(out_path)
    return out_path